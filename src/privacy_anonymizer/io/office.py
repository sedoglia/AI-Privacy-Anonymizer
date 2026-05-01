from __future__ import annotations

from pathlib import Path

from privacy_anonymizer.errors import MissingOptionalDependencyError
from privacy_anonymizer.io.base import FileAdapter, FileContent, WriteResult


class DocxAdapter(FileAdapter):
    extensions = {".docx"}

    def read_text(self, path: Path) -> FileContent:
        Document = _import_docx()
        document = Document(path)
        parts: list[str] = []
        parts.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        for section in document.sections:
            parts.extend(paragraph.text for paragraph in section.header.paragraphs if paragraph.text)
            parts.extend(paragraph.text for paragraph in section.footer.paragraphs if paragraph.text)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)
        return FileContent("\n".join(parts))

    def write_anonymized(
        self,
        source: Path,
        destination: Path,
        anonymized_text: str,
        keep_metadata: bool,
        replacements=None,
        original_text: str | None = None,
        source_content=None,
    ) -> WriteResult:
        del replacements, original_text, source_content
        Document = _import_docx()
        document = Document(source)
        _replace_docx_text(document, anonymized_text)
        metadata_stripped = False
        warnings = ["DOCX MVP: la sostituzione preserva il file ma può semplificare il testo in documenti con layout complesso."]
        if not keep_metadata:
            _strip_docx_metadata(document)
            try:
                accepted = _accept_track_changes(document)
                if accepted:
                    warnings.append(f"DOCX track-changes: {accepted} revisioni accettate o rimosse.")
            except ImportError:
                warnings.append("lxml non disponibile: track-changes non accettate esplicitamente.")
            metadata_stripped = True
        document.save(destination)
        return WriteResult(
            warnings=warnings,
            metadata_stripped=metadata_stripped,
        )


class XlsxAdapter(FileAdapter):
    extensions = {".xlsx"}

    def read_text(self, path: Path) -> FileContent:
        openpyxl = _import_openpyxl()
        workbook = openpyxl.load_workbook(path, data_only=False)
        values: list[str] = []
        for worksheet in workbook.worksheets:
            if worksheet.title:
                values.append(worksheet.title)
            for row in worksheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        values.append(cell.value)
                    if cell.comment and cell.comment.text:
                        values.append(cell.comment.text)
        return FileContent("\n".join(values))

    def write_anonymized(
        self,
        source: Path,
        destination: Path,
        anonymized_text: str,
        keep_metadata: bool,
        replacements=None,
        original_text: str | None = None,
        source_content=None,
    ) -> WriteResult:
        del replacements, original_text, source_content
        openpyxl = _import_openpyxl()
        workbook = openpyxl.load_workbook(source)
        replacements = _line_replacements(anonymized_text)
        index = 0
        for worksheet in workbook.worksheets:
            if worksheet.title and index < len(replacements):
                worksheet.title = _safe_xlsx_sheet_title(replacements[index], workbook.sheetnames, worksheet.title)
                index += 1
            for row in worksheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and index < len(replacements):
                        cell.value = replacements[index]
                        index += 1
                    if cell.comment and cell.comment.text and index < len(replacements):
                        cell.comment.text = replacements[index]
                        cell.comment.author = "Anonimo"
                        index += 1
        metadata_stripped = False
        if not keep_metadata:
            workbook.properties.creator = "Anonimo"
            workbook.properties.lastModifiedBy = "Anonimo"
            workbook.properties.company = None
            metadata_stripped = True
        workbook.save(destination)
        return WriteResult(metadata_stripped=metadata_stripped)


class PptxAdapter(FileAdapter):
    extensions = {".pptx"}

    def read_text(self, path: Path) -> FileContent:
        Presentation = _import_pptx()
        presentation = Presentation(path)
        values: list[str] = []
        for shape in _iter_pptx_shapes(presentation):
            if getattr(shape, "has_text_frame", False) and shape.text:
                values.append(shape.text)
        return FileContent("\n".join(values))

    def write_anonymized(
        self,
        source: Path,
        destination: Path,
        anonymized_text: str,
        keep_metadata: bool,
        replacements=None,
        original_text: str | None = None,
        source_content=None,
    ) -> WriteResult:
        del replacements, original_text, source_content
        Presentation = _import_pptx()
        presentation = Presentation(source)
        replacements = _line_replacements(anonymized_text)
        index = 0
        for shape in _iter_pptx_shapes(presentation):
            if getattr(shape, "has_text_frame", False) and shape.text and index < len(replacements):
                shape.text = replacements[index]
                index += 1
        metadata_stripped = False
        if not keep_metadata:
            core = presentation.core_properties
            core.author = "Anonimo"
            core.last_modified_by = "Anonimo"
            core.keywords = ""
            core.subject = ""
            core.comments = ""
            metadata_stripped = True
        presentation.save(destination)
        return WriteResult(metadata_stripped=metadata_stripped)


def _replace_docx_text(document, anonymized_text: str) -> None:
    replacements = _line_replacements(anonymized_text)
    index = 0
    for paragraph in document.paragraphs:
        if paragraph.text and index < len(replacements):
            paragraph.text = replacements[index]
            index += 1
    for section in document.sections:
        for paragraph in section.header.paragraphs:
            if paragraph.text and index < len(replacements):
                paragraph.text = replacements[index]
                index += 1
        for paragraph in section.footer.paragraphs:
            if paragraph.text and index < len(replacements):
                paragraph.text = replacements[index]
                index += 1
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if paragraph.text and index < len(replacements):
                        paragraph.text = replacements[index]
                        index += 1


def _strip_docx_metadata(document) -> None:
    core = document.core_properties
    core.author = "Anonimo"
    core.last_modified_by = "Anonimo"
    core.comments = ""
    core.keywords = ""
    core.subject = ""
    core.title = ""
    core.category = ""


W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _accept_track_changes(document) -> int:
    from lxml import etree

    accepted = 0
    body = document.element.body
    for ins in body.findall(f".//{{{W_NS}}}ins"):
        parent = ins.getparent()
        if parent is None:
            continue
        index = parent.index(ins)
        for child in list(ins):
            parent.insert(index, child)
            index += 1
        parent.remove(ins)
        accepted += 1
    for tag in ("del", "moveFrom"):
        for element in body.findall(f".//{{{W_NS}}}{tag}"):
            parent = element.getparent()
            if parent is not None:
                parent.remove(element)
                accepted += 1
    for tag in ("rsid", "rsidR", "rsidRPr", "rsidP", "rsidRDefault", "rsidTr"):
        for element in body.iter():
            attribute = f"{{{W_NS}}}{tag}"
            if attribute in element.attrib:
                del element.attrib[attribute]
    del etree
    return accepted


def _iter_pptx_shapes(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            yield shape
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                yield shape


def _line_replacements(text: str) -> list[str]:
    return text.splitlines()


def _safe_xlsx_sheet_title(candidate: str, existing_titles: list[str], current_title: str) -> str:
    forbidden = "[]:*?/\\"
    title = "".join("_" if char in forbidden else char for char in candidate).strip() or "Sheet"
    title = title[:31]
    taken = {title for title in existing_titles if title != current_title}
    if title not in taken:
        return title
    base = title[:28]
    counter = 1
    while f"{base}_{counter}" in taken:
        counter += 1
    return f"{base}_{counter}"[:31]


def _import_docx():
    try:
        from docx import Document
    except ImportError as exc:
        raise MissingOptionalDependencyError("python-docx", "office") from exc
    return Document


def _import_openpyxl():
    try:
        import openpyxl
    except ImportError as exc:
        raise MissingOptionalDependencyError("openpyxl", "office") from exc
    return openpyxl


def _import_pptx():
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise MissingOptionalDependencyError("python-pptx", "office") from exc
    return Presentation
